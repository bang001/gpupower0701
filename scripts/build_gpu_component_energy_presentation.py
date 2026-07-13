#!/usr/bin/env python3
"""Build the GPU component-energy methodology presentation.

The deck intentionally separates the 2026-07-08 RTX 3090 evidence snapshot
from the current active finalplan protocol. It describes GPU/device-level NVML
energy and effective microbenchmark coefficients, not physical circuit energy.
"""

from __future__ import annotations

import csv
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "gpu_component_energy_experiment_whitepaper_ko.pptx"
NOTES_PATH = OUT_DIR / "gpu_component_energy_experiment_whitepaper_ko.md"
PREVIEW_PATH = OUT_DIR / "gpu_component_energy_experiment_whitepaper_ko_contact_sheet.png"
PDF_PATH = OUT_DIR / "gpu_component_energy_experiment_whitepaper_ko.pdf"
RENDERED_PREVIEW_PATH = OUT_DIR / "gpu_component_energy_experiment_whitepaper_ko_rendered_contact_sheet.png"

W = 13.333
H = 7.5
FONT = "Malgun Gothic"
MONO = "Consolas"

INK = "172A3A"
TEXT = "263746"
MUTED = "667784"
PALE = "F2F5F6"
LINE = "D5DEE2"
WHITE = "FFFFFF"
TEAL = "187D77"
TEAL_PALE = "DFF1EF"
BLUE = "2E6F9E"
BLUE_PALE = "E4EFF7"
AMBER = "D78A26"
AMBER_PALE = "FBEED9"
RED = "B94B4B"
RED_PALE = "F8E5E5"
GREEN = "3A7D5D"
GREEN_PALE = "E4F1EA"
PURPLE = "6655A5"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def git_value(*args: str) -> str:
    return subprocess.check_output(["git", *args], cwd=ROOT, text=True).strip()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.04,
    name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs: Iterable[tuple[str, dict]], x, y, w, h, *, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    for value, options in runs:
        run = p.add_run()
        run.text = value
        run.font.name = options.get("font", FONT)
        run.font.size = Pt(options.get("size", 16))
        run.font.bold = options.get("bold", False)
        run.font.color.rgb = rgb(options.get("color", TEXT))
    return box


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill=WHITE,
    line=LINE,
    radius=True,
    text="",
    size=14,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.CENTER,
    name=None,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if text:
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = Inches(0.08)
        frame.margin_top = frame.margin_bottom = Inches(0.04)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width=1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x, y, w, h, *, fill=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    return shape


def add_status(slide, text, x, y, w, *, status="pass", size=11):
    colors = {
        "pass": (GREEN_PALE, GREEN),
        "accepted": (GREEN_PALE, GREEN),
        "warning": (AMBER_PALE, AMBER),
        "provisional": (AMBER_PALE, AMBER),
        "missing": (PALE, MUTED),
        "fail": (RED_PALE, RED),
        "reject": (RED_PALE, RED),
        "historical": (BLUE_PALE, BLUE),
    }
    fill, ink = colors.get(status, (PALE, MUTED))
    return add_box(slide, x, y, w, 0.32, fill=fill, line=fill, text=text, size=size, color=ink, bold=True)


def add_title(slide, title: str, kicker: str, slide_no: int, source: str = ""):
    add_text(slide, kicker.upper(), 0.55, 0.18, 7.0, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(slide, title, 0.55, 0.43, 12.0, 0.62, size=25, color=INK, bold=True)
    add_line(slide, 0.55, 1.08, 12.2, 1.08, color=LINE, width=1)
    if source:
        add_text(slide, source, 0.58, 7.12, 11.7, 0.2, size=7.2, color=MUTED, valign=MSO_ANCHOR.BOTTOM)
    add_text(slide, f"{slide_no:02d}", 12.35, 7.08, 0.4, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_section_label(slide, text, x, y, w, *, color=TEAL):
    add_text(slide, text, x, y, w, 0.3, size=10, color=color, bold=True)


def add_picture_cover(slide, path: Path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        crop_w = ih * box_ratio
        left = (iw - crop_w) / 2 / iw
        right = left
        top = bottom = 0
    else:
        crop_h = iw / box_ratio
        top = (ih - crop_h) / 2 / ih
        bottom = top
        left = right = 0
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.crop_left = left
    pic.crop_right = right
    pic.crop_top = top
    pic.crop_bottom = bottom
    return pic


def add_table(slide, rows: list[list[str]], x, y, w, h, *, widths=None, font_size=11, header_fill=INK, row_fills=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for col, width in zip(table.columns, widths):
            col.width = Inches(width)
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            fill = header_fill if ri == 0 else (row_fills[ri - 1] if row_fills else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            frame = cell.text_frame
            frame.clear()
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = frame.paragraphs[0]
            p.text = str(value)
            p.font.name = FONT
            p.font.size = Pt(font_size if ri else font_size - 0.3)
            p.font.bold = ri == 0
            p.font.color.rgb = rgb(WHITE if ri == 0 else TEXT)
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
    return table_shape


def read_strict_snapshot() -> list[dict[str, str]]:
    path = ROOT / "results" / "summary" / "rtx3090_strict_scope_fresh_ncu_component_coefficients_20260708.csv"
    with path.open(newline="") as f:
        return list(csv.DictReader(f))


def read_dram_reporting_policy() -> dict[str, str]:
    path = ROOT / "results" / "summary" / "rtx3090_dram_current_reporting_policy_20260712.csv"
    with path.open(newline="") as f:
        rows = list(csv.DictReader(f))
    if len(rows) != 1:
        raise ValueError(f"expected one DRAM reporting policy row in {path}")
    row = rows[0]
    if row["reporting_status"] != "provisional_reference_aligned_range" or row["strict_eligible"].lower() != "false":
        raise ValueError(f"invalid DRAM reporting policy status in {path}")
    return row


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    sha = git_value("rev-parse", "--short", "HEAD")
    branch = git_value("branch", "--show-current")
    dirty = bool(git_value("status", "--porcelain"))
    snapshot = read_strict_snapshot()

    # 1. Title
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_picture_cover(slide, ASSET_DIR / "system_measurement_architecture.png", 5.25, 0.0, 8.08, 7.5)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.45), Inches(7.5))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = rgb(WHITE); overlay.line.fill.background()
    add_text(slide, "GPU COMPONENT ENERGY", 0.68, 0.62, 4.8, 0.3, size=10, color=TEAL, bold=True)
    add_text(slide, "측정에서\nStrict Coefficient까지", 0.68, 1.10, 5.4, 1.7, size=31, color=INK, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, "NVML GPU/device-level energy · Treatment-Control · NCU path validation", 0.72, 3.04, 4.9, 0.7, size=15, color=TEXT, valign=MSO_ANCHOR.TOP)
    add_box(slide, 0.72, 4.18, 4.8, 1.25, fill=PALE, line=LINE, radius=False)
    add_text(slide, "핵심 해석", 0.95, 4.42, 1.2, 0.25, size=10, color=TEAL, bold=True)
    add_text(slide, "workload-dependent effective\nmicrobenchmark coefficient", 0.95, 4.72, 4.1, 0.52, size=18, color=INK, bold=True, valign=MSO_ANCHOR.TOP)
    state = "dirty worktree" if dirty else "clean worktree"
    add_text(slide, f"github.com/bang001/gpupower0701\nbranch {branch} · HEAD {sha} · {state} · 2026-07-12", 0.72, 6.45, 5.25, 0.55, size=9.5, color=MUTED, valign=MSO_ANCHOR.TOP)
    add_text(slide, "01", 12.35, 7.08, 0.4, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    # 2. Problem and objective
    slide = prs.slides.add_slide(blank); add_title(slide, "무엇을 해결하려는 실험인가", "PURPOSE", 2, "source: README.md; scripts/plan_platform_component_experiment.py")
    add_text(slide, "GPU 전체에서 관측한 에너지를, 경로가 검증된 미세벤치마크의 단위 작업당 증분으로 바꾼다.", 0.7, 1.28, 11.9, 0.52, size=19, color=INK, bold=True)
    stages = [
        ("GPU/device\nenergy", BLUE), ("Treatment−Control\nΔE", TEAL), ("NCU path\nvalidation", AMBER),
        ("Denominator\nvalidation", PURPLE), ("Power & stability\naudit", GREEN), ("Strict\ncoefficient", INK),
    ]
    x = 0.65
    for i, (label, color) in enumerate(stages):
        add_box(slide, x, 2.15, 1.72, 1.05, fill=WHITE, line=color, text=label, size=13, color=color, bold=True)
        if i < len(stages) - 1: add_arrow(slide, x + 1.79, 2.48, 0.31, 0.38, fill=LINE)
        x += 2.08
    questions = [
        "어떤 kernel과 control을 비교했는가?", "NVML은 어떤 scope를 측정했는가?", "NCU는 무엇을 검증했는가?",
        "분모는 expected인가 actual traffic인가?", "어떤 gate가 final 채택을 막는가?", "플랫폼별 완료 상태는 어디까지인가?",
    ]
    for i, q in enumerate(questions):
        col, row = i % 2, i // 2
        add_box(slide, 0.9 + col * 6.0, 3.85 + row * 0.72, 5.55, 0.52, fill=PALE, line=PALE, text=q, size=13, align=PP_ALIGN.LEFT)

    # 3. Interpretation boundary
    slide = prs.slides.add_slide(blank); add_title(slide, "측정하는 것과 측정하지 않는 것", "INTERPRETATION BOUNDARY", 3, "source: src/main.cu:755-768; scripts/analyze_matched_control_energy.py:1-10; NVIDIA NVML API")
    add_box(slide, 0.65, 1.35, 5.85, 4.85, fill=TEAL_PALE, line=TEAL, radius=False)
    add_section_label(slide, "말할 수 있는 것", 0.95, 1.62, 2.0, color=TEAL)
    add_text(slide, "NVML GPU/device-level energy를\nTreatment-Control 방식으로 차분하고,\nNCU counter로 경로·traffic denominator를 검증한\nworkload-dependent effective\nmicrobenchmark coefficient", 0.95, 2.02, 5.25, 1.72, size=15.5, color=INK, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, "• 특정 kernel · working set · occupancy · clock 상태에 의존\n• pJ/FLOP 또는 pJ/bit는 정규화 단위\n• GPU/device scope이며 외부 board meter 측정과 동일하지 않음", 1.0, 4.25, 5.0, 1.2, size=13, color=TEXT, valign=MSO_ANCHOR.TOP)
    add_box(slide, 6.85, 1.35, 5.85, 4.85, fill=RED_PALE, line=RED, radius=False)
    add_section_label(slide, "말하면 안 되는 것", 7.15, 1.62, 2.2, color=RED)
    forbidden = ["순수 Tensor Core 회로 에너지", "순수 register-file access energy", "순수 SRAM/HBM/GDDR bitcell energy", "모든 application의 GPU 고정 상수", "NCU가 직접 측정한 component energy"]
    for i, item in enumerate(forbidden):
        add_box(slide, 7.15, 2.10 + i * 0.70, 5.25, 0.5, fill=WHITE, line=RED_PALE, text="×  " + item, size=13, color=RED, bold=i == 4, align=PP_ALIGN.LEFT)
    add_text(slide, "board-level은 더 넓은 전원 경계를 암시한다. 본 코드가 명시적으로 기록하는 scope는 gpu_device_total_energy_counter다.", 0.9, 6.45, 11.7, 0.42, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # 4. System configuration, generated image 1
    slide = prs.slides.add_slide(blank); add_title(slide, "시스템 구성: 한 workload, 두 관측 채널", "SYSTEM", 4, "generated illustration + source: src/main.cu; src/nvml_energy.cpp; scripts/run_ncu_validation.sh")
    add_picture_cover(slide, ASSET_DIR / "system_measurement_architecture.png", 0.58, 1.22, 12.17, 5.45)
    add_box(slide, 0.85, 1.47, 2.2, 0.62, fill=WHITE, line=BLUE, text="Host CPU\nCUDA kernel dispatch", size=12, color=BLUE, bold=True)
    add_box(slide, 4.35, 1.47, 2.35, 0.62, fill=WHITE, line=INK, text="GPU accelerator\nbenchmark execution", size=12, color=INK, bold=True)
    add_box(slide, 8.3, 1.47, 2.08, 0.62, fill=WHITE, line=TEAL, text="NVML\ntotal-energy delta", size=12, color=TEAL, bold=True)
    add_box(slide, 10.35, 5.45, 2.05, 0.62, fill=WHITE, line=AMBER, text="NCU sidecar\npath counters", size=12, color=AMBER, bold=True)
    add_box(slide, 3.05, 6.0, 6.95, 0.5, fill=AMBER_PALE, line=AMBER, text="그림의 계측기는 개념 표현이다. 실제 코드는 외부 전력계를 사용하지 않고 NVML API를 호출한다.", size=10.5, color=TEXT)

    # 5. GPU hierarchy
    slide = prs.slides.add_slide(blank); add_title(slide, "GPU 계층과 실험이 겨냥하는 경로", "ARCHITECTURE", 5, "source: include/config.hpp; src/kernels.cu; docs/methodology/howitworks.md")
    add_box(slide, 0.7, 1.45, 3.4, 4.9, fill=PALE, line=INK, radius=False)
    add_text(slide, "SM", 0.95, 1.66, 0.7, 0.3, size=14, color=INK, bold=True)
    add_box(slide, 1.0, 2.15, 1.25, 0.75, fill=BLUE_PALE, line=BLUE, text="Tensor\nWMMA/HMMA", size=12, color=BLUE, bold=True)
    add_box(slide, 2.5, 2.15, 1.25, 0.75, fill=TEAL_PALE, line=TEAL, text="Register\noperands", size=12, color=TEAL, bold=True)
    add_box(slide, 1.0, 3.25, 2.75, 1.25, fill=WHITE, line=AMBER, text="Unified L1 / Shared physical resource\n서로 다른 instruction · address · bank/cache path", size=13, color=TEXT, bold=True)
    add_box(slide, 1.0, 4.88, 2.75, 0.72, fill=WHITE, line=LINE, text="Scheduler · scoreboard · control", size=12, color=MUTED)
    add_arrow(slide, 4.28, 3.2, 0.55, 0.48, fill=LINE)
    add_box(slide, 5.0, 2.15, 2.55, 2.9, fill=BLUE_PALE, line=BLUE, radius=False, text="L2 cache\n\nGPU-wide cache path\nld.global.cg로 L1 hit 억제", size=16, color=BLUE, bold=True)
    add_arrow(slide, 7.76, 3.2, 0.55, 0.48, fill=LINE)
    add_box(slide, 8.52, 2.15, 2.55, 2.9, fill=AMBER_PALE, line=AMBER, radius=False, text="External memory\n\nGDDR6X / HBM\ncontroller · PHY 포함 path", size=16, color=AMBER, bold=True)
    add_box(slide, 11.42, 1.63, 1.1, 0.58, fill=TEAL_PALE, line=TEAL, text="Shared", size=11, color=TEAL, bold=True)
    add_box(slide, 11.42, 2.42, 1.1, 0.58, fill=BLUE_PALE, line=BLUE, text="L1 hit", size=11, color=BLUE, bold=True)
    add_box(slide, 11.42, 3.21, 1.1, 0.58, fill=BLUE_PALE, line=BLUE, text="L2 hit", size=11, color=BLUE, bold=True)
    add_box(slide, 11.42, 4.0, 1.1, 0.58, fill=AMBER_PALE, line=AMBER, text="DRAM", size=11, color=AMBER, bold=True)
    add_text(slide, "같은 L1/Shared 하드웨어를 공유해도 접근 명령, arbitration, banking, cache lookup이 달라 effective path coefficient는 같을 필요가 없다.", 4.85, 5.55, 7.1, 0.62, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 6. Repository map
    slide = prs.slides.add_slide(blank); add_title(slide, "저장소 architecture map", "SOURCE OF TRUTH", 6, "source: src/; include/; scripts/; results/raw; results/ncu; results/summary")
    columns = [
        ("1 · CUDA/NVML 구현", ["src/main.cu", "src/kernels.cu", "src/nvml_energy.cpp", "src/result_writer.cpp"], BLUE),
        ("2 · 실행 계획", ["plan_platform_…py", "run_component_…py", "run_ncu_validation.sh", "generated commands.sh"], TEAL),
        ("3 · 분석·gate", ["power API/state audit", "NCU acceptance", "matched-control", "reliability/strict audit"], AMBER),
        ("4 · provenance", ["raw energy CSV", "NCU summary CSV", "matched detail CSV", "strict/package audit"], GREEN),
    ]
    for i, (title, items, color) in enumerate(columns):
        x = 0.55 + i * 3.15
        add_box(slide, x, 1.55, 2.75, 4.85, fill=WHITE, line=color, radius=False)
        add_box(slide, x, 1.55, 2.75, 0.62, fill=color, line=color, radius=False, text=title, size=13, color=WHITE, bold=True)
        for j, item in enumerate(items):
            add_box(slide, x + 0.2, 2.5 + j * 0.78, 2.35, 0.52, fill=PALE, line=PALE, text=item, size=12, color=TEXT)
        if i < 3: add_arrow(slide, x + 2.82, 3.48, 0.28, 0.36, fill=LINE)
    add_text(slide, "우선순위: 코드·CSV column → raw/NCU artifact → audit 결과 → 설명 문서", 1.0, 6.58, 11.3, 0.34, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 7. Active pairs
    slide = prs.slides.add_slide(blank); add_title(slide, "Current active treatment-control pair", "METHOD A · FINALPLAN", 7, "source: scripts/analyze_matched_control_energy.py:25-81; scripts/build_strict_component_summary.py:34-62")
    pairs = [
        ["Component/path", "Treatment", "Control", "Final unit"],
        ["Tensor MMA increment", "reg_mma", "reg_operand_only", "pJ/FLOP"],
        ["Shared scalar path", "shared_scalar_load_only", "clocked_empty", "pJ/bit"],
        ["Global L1 hit path", "global_l1_load_only", "global_addr_only", "pJ/bit"],
        ["L2 CG hit path", "l2_cg_load_only", "global_addr_only", "pJ/bit"],
        ["DRAM streaming sanity", "dram_cg_load_only", "global_addr_only", "pJ/bit · sanity"],
    ]
    add_table(slide, pairs, 0.65, 1.35, 12.0, 3.25, widths=[2.45,3.65,3.25,2.65], font_size=11.5, row_fills=[WHITE,PALE,WHITE,PALE,WHITE])
    add_picture_cover(slide, ASSET_DIR / "treatment_control_method.png", 0.65, 4.85, 6.0, 1.75)
    add_box(slide, 6.95, 4.85, 5.7, 1.75, fill=PALE, line=LINE, radius=False)
    add_text(slide, "Control의 역할", 7.2, 5.06, 1.7, 0.25, size=11, color=TEAL, bold=True)
    add_text(slide, "공통 loop · index · scheduler · epilogue 비용을 최대한 보존하고 target operation만 제거한다.", 7.2, 5.38, 5.0, 0.65, size=15, color=INK, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, "주의: l2_load_only는 normal global load이므로 strict L2-only 증거로 자동 채택하지 않는다.", 7.2, 6.05, 5.0, 0.34, size=10.5, color=RED, bold=True)

    # 8. Core pipeline
    slide = prs.slides.add_slide(blank); add_title(slide, "설명용 핵심 파이프라인", "PIPELINE · CONCEPTUAL", 8, "source: generated command packages; scripts/audit_*.py")
    core = ["Preflight", "Energy Run", "Power API /\nPower-state Audit", "NCU Sidecar", "Path Acceptance", "Matched-Control", "Reliability Audit", "Strict Summary"]
    colors = [INK, BLUE, TEAL, AMBER, PURPLE, BLUE, GREEN, INK]
    for i, (label, color) in enumerate(zip(core, colors)):
        row, col = divmod(i, 4)
        x, y = 0.75 + col * 3.1, 1.55 + row * 2.2
        add_box(slide, x, y, 2.55, 1.05, fill=WHITE, line=color, text=label, size=14, color=color, bold=True)
        add_text(slide, f"{i+1}", x + 0.12, y + 0.1, 0.28, 0.24, size=9, color=MUTED, bold=True)
        if col < 3: add_arrow(slide, x + 2.65, y + 0.34, 0.28, 0.36, fill=LINE)
    add_arrow(slide, 11.65, 2.73, 0.45, 0.55, fill=LINE)
    add_arrow(slide, 0.98, 3.04, 0.45, 0.55, fill=LINE)
    add_box(slide, 1.0, 6.05, 11.25, 0.52, fill=AMBER_PALE, line=AMBER, text="NCU는 Energy Run과 분리된 sidecar다. 같은 좌표의 실행 경로와 traffic을 검증하지만 energy를 측정하지 않는다.", size=13, color=TEXT, bold=True)

    # 9. Actual command pipeline
    slide = prs.slides.add_slide(blank); add_title(slide, "실제 command-package 파이프라인", "PIPELINE · IMPLEMENTATION", 9, "source: scripts/plan_platform_component_experiment.py:430-1016")
    lanes = [
        ("준비", ["Preflight", "Pipeline self-tests", "Stale output 격리", "Schema/revision smoke"], BLUE),
        ("측정·검증", ["Energy sweeps", "Power API audit", "Power-state audit", "NCU sidecar", "Path acceptance"], TEAL),
        ("분석·채택", ["Matched-control", "Reliability audit", "Instability audit", "Strict summary build", "Strict summary audit"], AMBER),
        ("패키징", ["Result manifest", "Package audit", "Gap / readiness / dashboard"], GREEN),
    ]
    y = 1.38
    for lane, items, color in lanes:
        add_box(slide, 0.65, y, 1.55, 0.92, fill=color, line=color, text=lane, size=13, color=WHITE, bold=True)
        usable = 10.1
        gap = 0.12
        node_w = (usable - gap * (len(items) - 1)) / len(items)
        x = 2.42
        for i, item in enumerate(items):
            add_box(slide, x, y, node_w, 0.92, fill=WHITE, line=color, text=item, size=10.5, color=TEXT, bold=True)
            if i < len(items)-1: add_arrow(slide, x + node_w + 0.02, y + 0.29, gap - 0.02, 0.3, fill=LINE)
            x += node_w + gap
        y += 1.28
    add_text(slide, "어느 중간 단계가 실패해도 triage artifact는 남기고, strict evidence package가 불완전하면 shell은 non-zero로 종료한다.", 0.85, 6.65, 11.6, 0.36, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 10. Platform profiles and blocks/SM sweep
    slide = prs.slides.add_slide(blank); add_title(slide, "Platform profile과 blocks/SM sweep", "PLATFORM", 10, "source: include/config.hpp:18-135; scripts/plan_platform_component_experiment.py:28-125")
    profile_rows = [
        ["Profile", "SM", "Shared/SM", "L1+Shared/SM", "L2", "max blocks/SM", "strict B/SM"],
        ["RTX 3090", "82", "100 KiB", "128 KiB", "6 MiB", "16", "8"],
        ["V100 32GB", "80", "96 KiB", "128 KiB", "6 MiB", "32", "32"],
        ["A100", "108", "164 KiB", "192 KiB", "40 MiB", "32", "16"],
        ["H100", "132", "228 KiB", "256 KiB", "50 MiB", "32", "16"],
    ]
    add_table(slide, profile_rows, 0.48, 1.22, 12.35, 1.82, widths=[1.55,0.65,1.55,1.85,1.1,2.2,1.7], font_size=8.7, row_fills=[WHITE,PALE,WHITE,PALE])
    slide.shapes.add_picture(str(ASSET_DIR / "platform_blocks_per_sm_sweep.png"), Inches(0.72), Inches(3.18), width=Inches(11.9), height=Inches(3.55))
    add_text(slide, "V100은 B4/B16 민감도와 strict B32를 측정한다. 요청 B/SM은 실제 동시 residency가 아니므로 NCU occupancy/resource로 확인한다.", 0.78, 6.78, 11.8, 0.3, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 11. Platform-specific W_SM path sweep
    slide = prs.slides.add_slide(blank); add_title(slide, "Platform별 W_SM path sweep", "EXPERIMENT MATRIX", 11, "source: planner profiles; exact-coordinate NCU acceptance decides the path")
    slide.shapes.add_picture(str(ASSET_DIR / "platform_wsm_path_sweep.png"), Inches(0.62), Inches(1.18), width=Inches(12.1), height=Inches(5.65))

    # 12. Host sequence
    slide = prs.slides.add_slide(blank); add_title(slide, "CUDA host 실행 sequence", "RAW ENERGY RUN", 12, "source: src/main.cu:628-650, 919-1010")
    seq = ["ITER\ncalibration", "Idle baseline\nmeasurement", "Global input\nwarm-up", "CUDA\nsync", "NVML\nbefore", "Benchmark\nkernel", "CUDA\nsync", "NVML\nafter", "SMID histogram\ncheck", "Scale idle +\nwrite CSV"]
    for i, label in enumerate(seq):
        row, col = divmod(i, 5)
        x, y = 0.55 + col * 2.55, 1.5 + row * 2.15
        color = TEAL if i in (4,7) else BLUE if i in (0,5) else AMBER if i in (1,8,9) else INK
        add_box(slide, x, y, 2.05, 0.95, fill=WHITE, line=color, text=label, size=12.5, color=color, bold=True)
        add_text(slide, f"{i+1}", x + 0.08, y + 0.08, 0.22, 0.2, size=8, color=MUTED, bold=True)
        if col < 4: add_arrow(slide, x + 2.1, y + 0.31, 0.28, 0.32, fill=LINE)
    add_arrow(slide, 11.8, 2.73, 0.45, 0.55, fill=LINE)
    add_arrow(slide, 0.85, 3.02, 0.45, 0.55, fill=LINE)
    add_box(slide, 0.8, 5.9, 11.75, 0.72, fill=PALE, line=LINE, text="SMID gate: unique SM = active_SM, total blocks = active_SM × blocks/SM, 각 사용 SM에 동일한 resident block 수가 배치되어야 한다.", size=13, color=INK, bold=True)

    # 13. Raw NVML energy
    slide = prs.slides.add_slide(blank); add_title(slide, "NVML energy와 idle-baseline 계산", "RAW NUMERATOR", 13, "source: src/main.cu:417-425, 692-720, 950-978; NVIDIA nvmlDeviceGetTotalEnergyConsumption")
    add_line(slide, 1.0, 2.05, 12.1, 2.05, color=INK, width=2.2)
    for x, label, color in [(1.25,"NVML before",TEAL),(4.25,"kernel execution",BLUE),(9.4,"NVML after",TEAL)]:
        add_box(slide, x, 1.53, 2.15, 0.75, fill=WHITE, line=color, text=label, size=13, color=color, bold=True)
    add_text(slide, "ΔE_raw = (E_after_mJ − E_before_mJ) / 1000", 0.8, 2.75, 5.7, 0.6, size=20, color=INK, bold=True, font=MONO)
    add_text(slide, "E_idle,scaled = E_idle × t_kernel / t_idle", 0.8, 3.65, 5.7, 0.6, size=20, color=INK, bold=True, font=MONO)
    add_text(slide, "E_net = ΔE_raw − E_idle,scaled", 0.8, 4.55, 5.7, 0.6, size=22, color=TEAL, bold=True, font=MONO)
    add_box(slide, 7.15, 2.72, 5.05, 2.45, fill=TEAL_PALE, line=TEAL, radius=False)
    add_text(slide, "Final numerator gate", 7.45, 3.02, 2.0, 0.3, size=11, color=TEAL, bold=True)
    gates = ["nvml_total_energy_supported = true", "energy_source = nvml_total_energy", "integration = total_energy_mj_delta", "scope = gpu_device_total_energy_counter"]
    for i, gate in enumerate(gates): add_text(slide, "✓  " + gate, 7.45, 3.42 + i * 0.42, 4.45, 0.3, size=12, color=INK, font=MONO)
    add_box(slide, 7.15, 5.42, 5.05, 1.02, fill=AMBER_PALE, line=AMBER, text="GetPowerUsage = fallback/provisional\npower.draw.* = metadata/diagnostic\nHopper module power · GPU memory power = component numerator 사용 금지", size=10.5, color=TEXT, bold=True)

    # 14. Memory matched control plus DRAM exception
    slide = prs.slides.add_slide(blank); add_title(slide, "Memory 차분: duration scaling과 DRAM 예외", "MATCHED CONTROL", 14, "source: scripts/analyze_matched_control_energy.py:620-658; scripts/run_component_regression_sweep.py")
    add_box(slide, 0.65, 1.35, 5.9, 4.95, fill=BLUE_PALE, line=BLUE, radius=False)
    add_section_label(slide, "Shared · Global L1 · L2", 0.98, 1.65, 2.7, color=BLUE)
    add_text(slide, "Treatment와 control의\nITER·elapsed가 다를 수 있다.", 0.98, 2.0, 5.0, 0.55, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "P_control =\nE_control,net / t_control", 1.12, 2.65, 4.7, 0.68, size=15.5, color=TEXT, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "E_control,scaled =\nP_control × t_treatment", 1.12, 3.45, 4.7, 0.68, size=15.5, color=TEXT, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "ΔE = E_treatment,net\n− E_control,scaled", 1.12, 4.28, 4.7, 0.68, size=16, color=BLUE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "elapsed ratio > 1.35 또는 신호가 너무 작으면 reject", 0.98, 5.22, 5.0, 0.4, size=12, color=RED, bold=True)
    add_box(slide, 6.85, 1.35, 5.9, 4.95, fill=AMBER_PALE, line=AMBER, radius=False)
    add_section_label(slide, "DRAM final pair", 7.18, 1.65, 2.1, color=AMBER)
    add_text(slide, "현행 DRAM sanity는\ntreatment/control work를 동일하게 잠근다.", 7.18, 2.0, 5.0, 0.58, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ITER_dram = ITER_addr", 7.18, 2.88, 4.8, 0.48, size=19, color=TEXT, bold=True, font=MONO)
    add_text(slide, "ΔE_DRAM = E_dram,net − E_addr,net", 7.18, 3.65, 5.0, 0.55, size=19, color=AMBER, bold=True, font=MONO)
    add_text(slide, "pair_energy_basis = matched_iters_net_energy\niter_ratio = 1", 7.18, 4.62, 4.9, 0.72, size=14, color=TEXT, bold=True, font=MONO, valign=MSO_ANCHOR.TOP)
    add_text(slide, "수정: ‘모든 memory pair가 duration-scaled’라는\n설명은 현재 코드와 맞지 않는다.", 7.18, 5.38, 4.9, 0.55, size=10.5, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 15. Tensor matched ITER
    slide = prs.slides.add_slide(blank); add_title(slide, "Tensor matched-ITER 예외", "TENSOR DIFFERENTIAL", 15, "source: scripts/run_component_regression_sweep.py:323-430; scripts/analyze_matched_control_energy.py:626-685")
    add_picture_cover(slide, ASSET_DIR / "treatment_control_method.png", 0.58, 1.25, 12.17, 5.62)
    add_box(slide, 0.88, 1.55, 3.15, 1.12, fill=WHITE, line=BLUE, text="CONTROL\nreg_operand_only", size=15, color=BLUE, bold=True)
    add_box(slide, 0.88, 4.55, 3.15, 1.12, fill=WHITE, line=TEAL, text="TREATMENT\nreg_mma", size=15, color=TEAL, bold=True)
    add_box(slide, 4.9, 1.65, 3.25, 1.05, fill=WHITE, line=INK, text="same W_SM · B/SM · active_SM · RF", size=13, color=INK, bold=True)
    add_box(slide, 4.9, 4.45, 3.25, 1.05, fill=WHITE, line=INK, text="same resolved ITER\nmax(treatment, control floor)", size=13, color=INK, bold=True)
    add_box(slide, 8.95, 2.85, 3.1, 1.65, fill=WHITE, line=AMBER, text="ΔE_tensor =\nE_reg_mma,net\n- E_reg_operand_only,net", size=13.5, color=INK, bold=True)
    add_status(slide, "pair_energy_basis = matched_iters_net_energy", 3.75, 6.16, 3.55, status="accepted", size=9.5)
    add_status(slide, "numerator_ITER = control_ITER", 7.48, 6.16, 2.6, status="accepted", size=9.5)
    add_status(slide, "iter_ratio = 1", 10.25, 6.16, 1.65, status="accepted", size=9.5)

    # 16. Energy vs NCU
    slide = prs.slides.add_slide(blank); add_title(slide, "Energy Run과 NCU Sidecar의 역할 분리", "TWO EVIDENCE CHANNELS", 16, "source: generated commands.sh comments #5 and #8; scripts/run_ncu_validation.sh")
    add_box(slide, 0.7, 1.45, 5.7, 4.75, fill=TEAL_PALE, line=TEAL, radius=False)
    add_section_label(slide, "Energy Run", 1.0, 1.78, 1.4, color=TEAL)
    energy_items = ["NCU detached", "약 10 s로 ITER calibration", "NVML total-energy before/after", "idle baseline 제거", "treatment-control ΔE 생성"]
    for i, item in enumerate(energy_items): add_box(slide, 1.0, 2.25+i*0.62, 5.0, 0.42, fill=WHITE, line=WHITE, text=item, size=12.5, color=TEXT, align=PP_ALIGN.LEFT)
    add_text(slide, "출력: Joule numerator", 1.0, 5.55, 4.95, 0.35, size=15, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 6.93, 1.45, 5.7, 4.75, fill=AMBER_PALE, line=AMBER, radius=False)
    add_section_label(slide, "NCU Sidecar", 7.23, 1.78, 1.7, color=AMBER)
    ncu_items = ["별도 profiler replay", "같은 W_SM/B/SM/RF/LR 좌표", "HMMA · cache hit/miss · bytes/access", "local/spill · occupancy · stall", "path acceptance와 denominator scale"]
    for i, item in enumerate(ncu_items): add_box(slide, 7.23, 2.25+i*0.62, 5.0, 0.42, fill=WHITE, line=WHITE, text=item, size=12.5, color=TEXT, align=PP_ALIGN.LEFT)
    add_text(slide, "출력: path/traffic evidence", 7.23, 5.55, 4.95, 0.35, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 2.05, 6.43, 9.25, 0.45, fill=RED_PALE, line=RED, text="NCU counter는 component energy를 직접 측정하지 않는다.", size=14, color=RED, bold=True)

    # 17. Tensor denominator
    slide = prs.slides.add_slide(blank); add_title(slide, "Tensor logical FLOP denominator", "DENOMINATOR · TENSOR", 17, "source: include/config.hpp:12-16; src/main.cu:781-791; scripts/analyze_ncu_path_acceptance.py")
    tensor_nodes = [
        ("active_SM", "실제 참여 SM"), ("blocks/SM", "SM당 block"), ("ITER", "반복 횟수"), ("reuse factor", "ITER당 MMA"),
    ]
    x = 0.65
    for i,(a,b) in enumerate(tensor_nodes):
        add_box(slide, x, 1.55, 2.25, 0.95, fill=WHITE, line=BLUE, text=f"{a}\n{b}", size=12.5, color=BLUE, bold=True)
        if i<3: add_text(slide, "×", x+2.28, 1.78, 0.35, 0.3, size=20, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        x += 2.65
    add_arrow(slide, 11.15, 1.78, 0.5, 0.42, fill=LINE)
    add_box(slide, 11.65, 1.55, 1.15, 0.95, fill=BLUE_PALE, line=BLUE, text="N_MMA", size=12, color=BLUE, bold=True)
    add_box(slide, 1.05, 3.0, 5.15, 1.25, fill=PALE, line=LINE, text="logical_FLOP = N_MMA × 8192\nwarp m16n16k16: 4096 FMA = 8192 FLOP", size=18, color=INK, bold=True)
    add_arrow(slide, 6.35, 3.38, 0.52, 0.42, fill=LINE)
    add_box(slide, 7.05, 3.0, 5.15, 1.25, fill=TEAL_PALE, line=TEAL, text="pJ/FLOP = ΔE_tensor × 10¹²\n÷ logical_FLOP", size=19, color=TEAL, bold=True)
    add_box(slide, 0.95, 5.1, 11.45, 1.1, fill=AMBER_PALE, line=AMBER, text="NCU의 HMMA instruction은 denominator 자체가 아니다. Treatment에 Tensor instruction이 있고, control에는 없으며, spill/local 및 과도한 memory traffic이 없는지를 검증한다.", size=14, color=TEXT, bold=True)

    # 18. Memory denominator
    slide = prs.slides.add_slide(blank); add_title(slide, "Memory traffic denominator와 NCU scale", "DENOMINATOR · MEMORY", 18, "source: scripts/analyze_matched_control_energy.py:339-399, 643-658")
    add_box(slide, 0.65, 1.45, 3.4, 1.08, fill=BLUE_PALE, line=BLUE, text="Expected bytes (energy run)\nactive_SM × B/SM × ITER × LR × 1024", size=14, color=BLUE, bold=True)
    add_arrow(slide, 4.2, 1.78, 0.5, 0.42, fill=LINE)
    add_box(slide, 4.85, 1.45, 3.4, 1.08, fill=AMBER_PALE, line=AMBER, text="NCU scale\nactual bytes / expected bytes", size=14, color=AMBER, bold=True)
    add_arrow(slide, 8.4, 1.78, 0.5, 0.42, fill=LINE)
    add_box(slide, 9.05, 1.45, 3.65, 1.08, fill=TEAL_PALE, line=TEAL, text="Final denominator\nenergy expected × NCU scale", size=14, color=TEAL, bold=True)
    add_text(slide, "pJ/byte = ΔE × 10¹² / final_bytes", 1.0, 3.05, 5.5, 0.55, size=20, color=INK, bold=True, font=MONO)
    add_text(slide, "pJ/bit = pJ/byte ÷ 8", 7.15, 3.05, 4.8, 0.55, size=20, color=INK, bold=True, font=MONO)
    provenance = [
        ("ncu_actual_exact", "mode·W·B·SM·factor exact match", "strict final 허용", GREEN_PALE, GREEN),
        ("ncu_actual_same_working_set", "같은 W/B/SM, factor scale 재사용", "민감도·보조", AMBER_PALE, AMBER),
        ("expected_no_ncu_match", "일치 NCU row 없음", "strict memory reject", RED_PALE, RED),
    ]
    for i,(name, meaning, use, fill, color) in enumerate(provenance):
        y=4.25+i*0.72
        add_box(slide, 0.9, y, 3.15, 0.52, fill=fill, line=color, text=name, size=11, color=color, bold=True)
        add_text(slide, meaning, 4.3, y, 4.55, 0.52, size=12.5, color=TEXT)
        add_status(slide, use, 9.35, y+0.1, 2.7, status="accepted" if i==0 else "warning" if i==1 else "fail", size=10)

    # 19. NCU validation
    slide = prs.slides.add_slide(blank); add_title(slide, "NCU path acceptance와 spill evidence", "VALIDATION", 19, "generated illustration + source: scripts/analyze_ncu_path_acceptance.py; scripts/run_ncu_validation.sh")
    add_picture_cover(slide, ASSET_DIR / "ncu_audit_validation.png", 0.58, 1.2, 12.17, 5.75)
    labels = [
        ("Tensor", "HMMA > 0\ncontrol HMMA = 0", 0.75, 1.35, BLUE),
        ("Shared", "shared bytes/access\nlow global traffic", 0.75, 2.45, TEAL),
        ("Global L1", "path hit ≥ 95%\nlow L2/DRAM ratio", 0.75, 3.55, BLUE),
        ("L2 CG", "L1 hit ≤ 1%\nL2 read hit ≥ 95%", 0.75, 4.65, PURPLE),
        ("Spill/Stall", "local bytes = 0\nlong scoreboard context", 0.75, 5.75, AMBER),
    ]
    for title, body, x,y,color in labels:
        add_box(slide, x,y,2.4,0.82,fill=WHITE,line=color,text=f"{title}\n{body}",size=10.5,color=color,bold=True)
    add_box(slide, 9.7, 2.2, 2.4, 1.15, fill=GREEN_PALE, line=GREEN, text="accepted\npath + denominator evidence", size=13, color=GREEN, bold=True)
    add_box(slide, 9.7, 4.4, 2.4, 1.15, fill=RED_PALE, line=RED, text="rejected/provisional\nmissing or contaminated", size=13, color=RED, bold=True)
    add_box(slide, 3.25, 6.18, 7.35, 0.58, fill=WHITE, line=LINE, text="Hit rate만으로는 부족하다: access count · bytes · hit/miss sectors · local traffic · stall/status를 함께 보존한다.", size=10.5, color=INK, bold=True)

    # 20. Audit funnel / states
    slide = prs.slides.add_slide(blank); add_title(slide, "Audit 단계별 상태 체계", "EVIDENCE FUNNEL", 20, "source: scripts/audit_power_api_measurements.py; analyze_ncu_path_acceptance.py; audit_component_reliability.py; audit_platform_result_package.py")
    state_rows = [
        ("Power API audit", ["final_candidate", "provisional", "reject"], TEAL),
        ("NCU path acceptance", ["accepted", "provisional", "rejected"], AMBER),
        ("Reliability audit", ["accepted", "accepted_with_caution", "accepted_low_stability", "accepted_sanity", "reject"], BLUE),
        ("Package / goal readiness", ["pass", "warning", "missing", "fail"], GREEN),
    ]
    y=1.45
    for idx,(stage,states,color) in enumerate(state_rows):
        add_box(slide, 0.7,y,2.35,0.88,fill=color,line=color,text=stage,size=13,color=WHITE,bold=True)
        x=3.35
        available=9.15
        gap=0.10
        sw=(available-gap*(len(states)-1))/len(states)
        for s in states:
            status="pass" if s in {"accepted","final_candidate","pass"} else "warning" if "caution" in s or "low" in s or s in {"provisional","accepted_sanity","warning"} else "missing" if s=="missing" else "fail"
            add_status(slide,s,x,y+0.28,sw,status=status,size=8.7 if len(states)>4 else 10)
            x+=sw+gap
        if idx<3: add_arrow(slide, 1.62, y+0.95, 0.5, 0.38, fill=LINE)
        y+=1.28
    add_box(slide, 1.1, 6.55, 11.1, 0.42, fill=RED_PALE, line=RED, text="서로 다른 단계의 상태를 하나의 ‘confidence’ 등급처럼 합치지 않는다.", size=13, color=RED, bold=True)

    # 21. RTX snapshot results and current DRAM reporting policy
    slide = prs.slides.add_slide(blank); add_title(slide, "RTX 3090 snapshot과 DRAM 최신 보고 범위", "RESULTS · EVIDENCE BOUNDARY", 21, "source: strict snapshot; current-protocol reaudit; rtx3090_dram_current_reporting_policy_20260712.csv")
    strict = {row["component"]: row for row in snapshot}
    dram = read_dram_reporting_policy()
    tensor_value = float(strict["Tensor MMA incremental"]["median"])
    shared_value = float(strict["Shared scalar path"]["median"])
    l1_value = float(strict["Global L1 hit path"]["median"])
    l2_value = float(strict["L2 CG hit path"]["median"])
    dram_low = float(dram["estimate_low"])
    dram_high = float(dram["estimate_high"])
    add_status(slide, "2026-07-08 accepted strict snapshot: 4 rows", 0.68, 1.18, 3.25, status="historical", size=9.6)
    add_status(slide, "2026-07-12 current-protocol reaudit: fail", 4.08, 1.18, 3.15, status="fail", size=9.6)
    add_box(slide, 7.42, 1.18, 5.15, 0.32, fill=PALE, line=LINE, text="all rows: nvml_total_energy · total_energy_mj_delta · gpu_device_total_energy_counter", size=7.7, color=MUTED)
    add_text(slide, "Tensor · pJ/FLOP", 0.75, 1.7, 2.5, 0.25, size=10.5, color=BLUE, bold=True)
    add_box(slide, 0.75, 2.02, 4.55, 0.48, fill=BLUE_PALE, line=BLUE, radius=False)
    add_box(slide, 0.75, 2.02, 2.9, 0.48, fill=BLUE, line=BLUE, radius=False)
    add_text(slide, f"{tensor_value:.3f}", 3.85, 2.02, 1.25, 0.48, size=17, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "Historical on-chip snapshot · pJ/bit", 5.85, 1.7, 3.4, 0.25, size=10.5, color=TEAL, bold=True)
    mem = [("Shared",shared_value,TEAL),("Global L1",l1_value,BLUE),("L2 CG",l2_value,PURPLE)]
    maxv=1.2
    for i,(name,val,color) in enumerate(mem):
        y=1.98+i*0.25
        add_text(slide,name,5.85,y,1.0,0.2,size=7.8,color=TEXT,bold=True)
        add_box(slide,6.9,y,2.25,0.18,fill=PALE,line=PALE,radius=False)
        add_box(slide,6.9,y,max(0.08,2.25*val/maxv),0.18,fill=color,line=color,radius=False)
        add_text(slide,f"{val:.3f}",9.24,y-0.03,0.62,0.24,size=8.3,color=color,bold=True,align=PP_ALIGN.RIGHT)
    add_box(slide, 10.0, 1.72, 2.35, 0.92, fill=AMBER_PALE, line=AMBER, radius=False)
    add_text(slide, "DRAM provisional band", 10.18, 1.82, 2.0, 0.22, size=9.2, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"{dram_low:.3f}-{dram_high:.3f}\npJ/bit", 10.18, 2.04, 2.0, 0.48, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    detail_rows = [
        ["Result", "Snapshot pair → active pair", "Condition", "Denominator / NCU evidence", "Reliability / current"],
        [f"Tensor\n{tensor_value:.3f} pJ/FLOP", "reg_mma - reg_operand_only\n→ same pair", "W2048 KiB · B16 · SM82\nRF8,16", "logical FLOP\nHMMA treatment · legacy control/spill schema", "accepted · medium\ncurrent rerun"],
        [f"Shared\n{shared_value:.3f} pJ/bit", "shared_scalar - clocked_empty\n→ same pair", "W64 KiB · B16 · SM82\nLR8", "ncu_actual_exact\nshared bytes 1.0748e12", "accepted · medium\ncurrent rerun"],
        [f"Global L1\n{l1_value:.3f} pJ/bit", "global_l1 - clocked_empty\n→ global_l1 - global_addr", "W16 KiB · B16 · SM82\nLR4", "ncu_actual_exact\nL1 hit 99.9995%", "accepted · medium\npair changed"],
        [f"L2 CG\n{l2_value:.3f} pJ/bit", "l2_cg - clocked_empty\n→ l2_cg - global_addr", "W64 KiB · B16 · SM82\nLR4,8", "ncu_actual_exact\nL1 ≈0% · L2 hit ≈99.985%", "accepted · medium\npair changed"],
        [f"DRAM cumulative path\n{dram_low:.3f}-{dram_high:.3f} pJ/bit", "target: dram_cg - global_addr\nlegacy clocked_empty excluded", "W8192 KiB · matched ITER\ncurrent raw pair 없음", "target denominator\nexact NCU dram_bytes × 8", "provisional range\nstrict eligible = false"],
    ]
    add_table(slide, detail_rows, 0.55, 3.05, 12.2, 3.15, widths=[1.45,3.0,2.35,3.15,2.25], font_size=7.5, row_fills=[WHITE,PALE,WHITE,PALE,WHITE])
    add_box(slide, 0.75, 6.42, 11.55, 0.45, fill=RED_PALE, line=RED, text="DRAM 26.709-28.409 pJ/bit는 최신 provisional cumulative-path band다. matched-ITER address-control raw pair가 없으므로 accepted 실측값이나 physical GDDR6X energy로 쓰지 않는다.", size=9.8, color=RED, bold=True)

    # 22. Cross-platform state and checklist
    slide = prs.slides.add_slide(blank); add_title(slide, "Cross-platform 상태 · 한계 · 재현 checklist", "STATUS & NEXT ACTION", 22, "source: component_energy_goal_readiness_audit_20260712.csv; platform package audits; official NVIDIA NVML docs")
    matrix = [
        ["Evidence stage", "RTX 3090", "V100", "A100", "H100"],
        ["Repository profile", "확인", "확인", "확인", "확인"],
        ["Command package", "확인", "확인", "확인", "확인"],
        ["Raw energy artifact", "historical 있음", "missing", "missing", "missing"],
        ["Fresh NCU acceptance", "current fail", "missing", "missing", "missing"],
        ["Reliability audit", "historical 있음", "missing", "missing", "missing"],
        ["Strict component summary", "historical 4-row", "missing", "missing", "missing"],
        ["Package audit complete", "아니오", "아니오", "아니오", "아니오"],
    ]
    add_table(slide,matrix,0.55,1.23,7.6,4.75,widths=[2.4,1.35,1.2,1.2,1.2],font_size=9.6,row_fills=[WHITE,PALE,WHITE,PALE,WHITE,PALE,WHITE])
    add_box(slide,8.45,1.23,4.3,2.15,fill=TEAL_PALE,line=TEAL,radius=False)
    add_text(slide,"재현 전 필수",8.75,1.48,1.6,0.28,size=11,color=TEAL,bold=True)
    add_text(slide,"1. strict preflight + runtime profile\n2. clean architecture-specific build\n3. total-energy scope row 확인\n4. exact-coordinate treatment/control NCU\n5. package audit 0 fail · 0 missing",8.75,1.88,3.55,1.2,size=12,color=TEXT,bold=True,valign=MSO_ANCHOR.TOP)
    add_box(slide,8.45,3.65,4.3,2.33,fill=AMBER_PALE,line=AMBER,radius=False)
    add_text(slide,"남은 해석 한계",8.75,3.9,1.7,0.28,size=11,color=AMBER,bold=True)
    add_text(slide,"• target profile은 SKU 보편값이 아님\n• H100은 WMMA, WGMMA/TMA가 아님\n• DRAM은 hierarchy sanity\n• coefficient는 workload-dependent\n• external board power가 아님",8.75,4.28,3.55,1.35,size=12,color=TEXT,bold=True,valign=MSO_ANCHOR.TOP)
    add_box(slide,0.78,6.25,11.9,0.55,fill=INK,line=INK,text="완료 조건: Energy + NCU + denominator + power/stability + strict/package audit가 같은 좌표와 provenance로 연결될 것",size=13,color=WHITE,bold=True)

    assert len(prs.slides) == 22
    return prs


SLIDE_NOTES = """# GPU Component Energy Experiment Presentation Evidence Notes

## Deck identity

- Repository: `https://github.com/bang001/gpupower0701`
- Generated from the current worktree on 2026-07-12.
- The title slide records the checked-out branch and HEAD SHA.
- The deck has 22 slides and three generated technical illustrations.

## Source-of-truth order

1. CUDA/NVML implementation and generated planner commands
2. Raw CSV columns and artifact provenance
3. NCU, power, reliability, strict, and package audit artifacts
4. Active methodology/platform documentation
5. Official NVIDIA NVML API documentation for API semantics

## Slide evidence

1. **Title** — repository/HEAD from Git; generated system illustration.
2. **Purpose** — `README.md`, `scripts/plan_platform_component_experiment.py`.
3. **Interpretation boundary** — `src/main.cu`, `scripts/analyze_matched_control_energy.py`.
4. **System** — `src/main.cu`, `src/nvml_energy.cpp`, `scripts/run_ncu_validation.sh`. The pictured instrument is conceptual; actual code uses NVML, not an external meter.
5. **Hierarchy** — `include/config.hpp`, `src/kernels.cu`, `docs/methodology/howitworks.md`.
6. **Repository map** — `src/`, `include/`, `scripts/`, `results/`.
7. **Active pairs** — `scripts/analyze_matched_control_energy.py:25-81`, `scripts/build_strict_component_summary.py:34-62`.
8. **Core pipeline** — generated command packages.
9. **Actual pipeline** — `scripts/plan_platform_component_experiment.py:430-1016`.
10. **Profiles and blocks/SM sweep** — `include/config.hpp:18-135`, planner profiles, and `platform_blocks_per_sm_sweep.png`. V100 uses B4/B16 sensitivity points and the strict B32 anchor; requested blocks/SM still needs NCU occupancy/resource validation.
11. **Platform W_SM path sweep** — planner profiles, generated `*_command_plan.md`, and `platform_wsm_path_sweep.png`. Shared is a separate address-space path; only global-memory candidates are interpreted across L1/L2/DRAM after exact-coordinate NCU acceptance.
12. **Host sequence** — `src/main.cu:628-650,919-1010`.
13. **Raw energy** — `src/main.cu:417-425,692-720,950-978`.
14. **Memory differential** — `scripts/analyze_matched_control_energy.py`. Shared/L1 use duration-scaled control power. Current L2 CG/DRAM CG finalplans require matched ITER and direct net-energy subtraction.
15. **Tensor differential** — `scripts/run_component_regression_sweep.py:323-430`; matched-ITER fields in matched detail.
16. **Energy vs NCU** — planner comments #5 and #8.
17. **Tensor denominator** — `include/config.hpp:12-16`, `src/main.cu:781-791`. NCU HMMA is validation evidence, not the final FLOP denominator.
18. **Memory denominator** — `scripts/analyze_matched_control_energy.py:339-399,643-658`. `expected_no_ncu_match` is not eligible for strict memory coefficients.
19. **NCU acceptance** — `scripts/analyze_ncu_path_acceptance.py`. Acceptance uses path-specific rates plus access/byte/local/stall evidence, not aggregate hit rate alone.
20. **Audit states** — each audit script owns a distinct state vocabulary; states are not collapsed into one grade.
21. **RTX 3090 snapshot and DRAM policy** — the historical strict snapshot contains four rows and fails the active control/schema reaudit. DRAM is now shown as the 26.709-28.409 pJ/bit `provisional_reference_aligned_range` from `rtx3090_dram_current_reporting_policy_20260712.csv`; no completed matched-ITER address-control raw pair exists, so the band is not strict measured evidence.
22. **Cross-platform state** — current readiness/package audits. A command package is not evidence that a target-node experiment completed.

## Corrected statements

- The strict RTX 3090 snapshot has four historical components, not five. DRAM is displayed separately as a provisional cumulative-path reporting band.
- Global L1 and L2 historical rows used `clocked_empty`; active finalplan uses `global_addr_only`.
- Current L2 CG and DRAM CG finalplans use pair-locked identical ITER; they are not duration-scaled.
- `l2_load_only` is a normal global-load capacity diagnostic, not automatic strict L2-only evidence.
- The H100 implementation is FP16 WMMA compatibility code, not Hopper-native WGMMA/TMA.
- NVML total-energy delta is labeled GPU/device scope. It is not represented as an external whole-board meter.

## Official NVIDIA references

- NVML API Reference, `nvmlDeviceGetTotalEnergyConsumption`: https://docs.nvidia.com/deploy/nvml-api/group__nvmlDeviceQueries.html
  - Reports total GPU energy in mJ since driver reload; supported on Volta or newer fully supported devices.
- NVML API Reference, `nvmlDeviceGetPowerUsage`: https://docs.nvidia.com/deploy/nvml-api/group__nvmlDeviceQueries.html
  - Ampere except GA100 and newer: 1-second averaged power.
  - GA100 and older: instantaneous power.
  - The API reports GPU power and associated circuitry such as memory.

## Generated illustrations

The three raster illustrations were generated with the built-in image generation tool using the `scientific-educational` use case. Exact labels, equations, and status values are PowerPoint text/shapes so that technical claims are not dependent on generated in-image text.
"""


def render_contact_sheet(prs: Presentation) -> None:
    """Render an approximate contact sheet for layout inspection.

    This is not a PowerPoint rendering engine. It visualizes shape bounds, text,
    fills, and images closely enough to catch clipped or crowded layouts.
    """

    sw, sh = 1066, 600
    scale_x, scale_y = sw / W, sh / H
    font_path = "/usr/share/fonts/truetype/unfonts-core/UnDotum.ttf"
    bold_path = "/usr/share/fonts/truetype/unfonts-core/UnDotumBold.ttf"
    thumbs: list[Image.Image] = []
    for slide in prs.slides:
        canvas = Image.new("RGB", (sw, sh), "white")
        draw = ImageDraw.Draw(canvas)
        for shape in slide.shapes:
            x = int(shape.left / 914400 * scale_x)
            y = int(shape.top / 914400 * scale_y)
            w = max(1, int(shape.width / 914400 * scale_x))
            h = max(1, int(shape.height / 914400 * scale_y))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    from io import BytesIO
                    im = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                    im.thumbnail((w, h))
                    canvas.paste(im, (x + (w - im.width)//2, y + (h - im.height)//2))
                except Exception:
                    draw.rectangle((x,y,x+w,y+h), outline="#D5DEE2")
                continue
            fill = "#FFFFFF"
            try:
                if shape.fill.type is not None and shape.fill.fore_color.rgb:
                    fill = f"#{shape.fill.fore_color.rgb}"
            except Exception:
                pass
            if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.TABLE}:
                draw.rounded_rectangle((x,y,x+w,y+h), radius=4, fill=fill, outline="#D5DEE2", width=1)
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text = shape.text.strip()
                try:
                    p = shape.text_frame.paragraphs[0]
                    size = p.runs[0].font.size.pt if p.runs and p.runs[0].font.size else 12
                    bold = bool(p.runs and p.runs[0].font.bold)
                    color = "#263746"
                    if p.runs and p.runs[0].font.color.type is not None and p.runs[0].font.color.rgb:
                        color = f"#{p.runs[0].font.color.rgb}"
                except Exception:
                    size, bold, color = 12, False, "#263746"
                px = max(6, int(size * 1.05))
                path = bold_path if bold and Path(bold_path).exists() else font_path
                font = ImageFont.truetype(path, px)
                lines = text.splitlines()
                line_h = px + 2
                max_lines = max(1, h // line_h)
                lines = lines[:max_lines]
                ty = y + max(2, (h - len(lines)*line_h)//2)
                for line in lines:
                    draw.text((x+4,ty), line[:70], font=font, fill=color)
                    ty += line_h
        thumbs.append(canvas.resize((400,225)))
    sheet = Image.new("RGB", (1600, 1350), "#E9EEF0")
    for i, thumb in enumerate(thumbs):
        col, row = i % 4, i // 4
        sheet.paste(thumb, (col*400, row*225))
    sheet.save(PREVIEW_PATH)


def validate(prs: Presentation) -> None:
    assert len(prs.slides) == 22
    for index, slide in enumerate(prs.slides, 1):
        assert slide.shapes, f"slide {index} is empty"
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0, f"negative placement on slide {index}"
            assert shape.left + shape.width <= prs.slide_width + 2, f"shape past right edge on slide {index}: {shape.name}"
            assert shape.top + shape.height <= prs.slide_height + 2, f"shape past bottom edge on slide {index}: {shape.name}"
    assert len(read_strict_snapshot()) == 4


def export_rendered_review() -> None:
    """Export a PDF and rendered contact sheet when LibreOffice is available."""

    libreoffice = shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    montage = shutil.which("montage")
    if not libreoffice:
        print("skip PDF render: libreoffice not found")
        return
    render_dir = Path("/tmp/gpupower_component_energy_presentation")
    if render_dir.exists():
        shutil.rmtree(render_dir)
    render_dir.mkdir(parents=True)
    subprocess.run(
        [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", str(render_dir), str(PPTX_PATH)],
        check=True,
    )
    generated_pdf = render_dir / PPTX_PATH.with_suffix(".pdf").name
    shutil.copy2(generated_pdf, PDF_PATH)
    if not (pdftoppm and montage):
        print("skip rendered contact sheet: pdftoppm or montage not found")
        return
    png_dir = render_dir / "png"
    png_dir.mkdir()
    subprocess.run([pdftoppm, "-png", "-r", "100", str(PDF_PATH), str(png_dir / "slide")], check=True)
    slide_pngs = sorted(png_dir.glob("slide-*.png"))
    subprocess.run(
        [montage, *(str(path) for path in slide_pngs), "-thumbnail", "400x225", "-tile", "4x6", "-geometry", "+0+0", str(RENDERED_PREVIEW_PATH)],
        check=True,
    )


def main() -> int:
    subprocess.run([sys.executable, str(ROOT / "scripts" / "plot_platform_sweep_design.py")], check=True)
    subprocess.run([sys.executable, str(ROOT / "scripts" / "plot_dram_reporting_policy.py")], check=True)
    for path in [
        ASSET_DIR / "system_measurement_architecture.png",
        ASSET_DIR / "treatment_control_method.png",
        ASSET_DIR / "ncu_audit_validation.png",
        ASSET_DIR / "platform_blocks_per_sm_sweep.png",
        ASSET_DIR / "platform_wsm_path_sweep.png",
        ASSET_DIR / "platform_capacity_context.png",
    ]:
        if not path.exists():
            raise FileNotFoundError(path)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    validate(prs)
    prs.save(PPTX_PATH)
    NOTES_PATH.write_text(SLIDE_NOTES, encoding="utf-8")
    render_contact_sheet(prs)
    export_rendered_review()
    print(f"wrote {PPTX_PATH}")
    print(f"wrote {NOTES_PATH}")
    print(f"wrote {PREVIEW_PATH}")
    if PDF_PATH.exists():
        print(f"wrote {PDF_PATH}")
    if RENDERED_PREVIEW_PATH.exists():
        print(f"wrote {RENDERED_PREVIEW_PATH}")
    print(f"slides={len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
